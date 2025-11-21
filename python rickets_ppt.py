# rickets_ppt.py
# Author:  <you>
# Date:    2025-11-21
# Purpose: Detailed veterinary PowerPoint on Rickets in the Dog
# -------------------------------------------------------------

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os, requests, io

# ------------------------------------------------------------------
# Helper: download image once, return local path
# ------------------------------------------------------------------
def fetch_img(url, fname):
    if not os.path.exists(fname):
        r = requests.get(url, timeout=30)
        with open(fname, "wb") as f:
            f.write(r.content)
    return fname

# Radiology images (Wikimedia Commons – public domain)
rad_urls = {
    "rickets_lateral":"https://upload.wikimedia.org/wikipedia/commons/3/3e/Rickets_lateral.jpg",
    "rickets_ap":"https://upload.wikimedia.org/wikipedia/commons/2/2b/Rickets_AP.jpg"
}
for k, v in rad_urls.items():
    fetch_img(v, f"{k}.jpg")

# ------------------------------------------------------------------
# Helper: add a slide with title + content (bullets or picture)
# ------------------------------------------------------------------
def add_slide(prs, title, bullets=None, picture=None, notes=None):
    layout = prs.slide_layouts[5]   # blank layout
    slide = prs.slides.add_slide(layout)

    # title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)

    if bullets:
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        tf = body.text_frame
        for line in bullets:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(255,255,255)
            p.level = 0
    if picture:
        slide.shapes.add_picture(picture, Inches(5.5), Inches(1.8), width=Inches(4))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

# ------------------------------------------------------------------
# RGB helper
# ------------------------------------------------------------------
from pptx.dml.color import RGBColor

# ------------------------------------------------------------------
# Build presentation
# ------------------------------------------------------------------
prs = Presentation()
# set widescreen
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# global background: dark blue
def dark_blue_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)

# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
dark_blue_bg(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Rickets in the Dog"
subtitle.text = "Etiology – Clinical Signs – Diagnosis – Treatment – Management\nVeterinary Continuing Education"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(221,221,221)

# ------------------------------------------------------------------
# 1. Learning objectives
add_slide(prs,
          "Learning Objectives",
          bullets=["Define rickets and explain the underlying pathophysiology",
                   "List the principal causes in puppies",
                   "Recognise clinical & radiological signs",
                   "Formulate a diagnostic plan",
                   "Outline evidence-based treatment and long-term management"],
          notes="Emphasise that rickets ONLY occurs in growing animals – the growth plate is the target.")

# ------------------------------------------------------------------
# 2. What is Rickets?
add_slide(prs,
          "What is Rickets?",
          bullets=["Metabolic bone disease of young, growing dogs",
                   "Failure of mineralisation at the zone of provisional calcification → widened, irregular physes",
                   "Consequence of deficient Ca, P or vitamin D",
                   "End-result: soft osteoid → bowed limbs, fractures, pain"],
          notes="Pathology is failure of vascular invasion + mineralisation in metaphysis.")

# ------------------------------------------------------------------
# 3. Etiology & Risk Factors
add_slide(prs,
          "Etiology & Risk Factors",
          bullets=["Nutritional (90 % of cases)",
                   "  – All-meat diet → low Ca, low vitamin D",
                   "  – Incorrect raw/cooked homemade diet",
                   "  – Excessive Ca (≥3× normal) → secondary rickets-like syndrome in giant breeds",
                   "Lack of sunlight → ↓cutaneous vitamin D3",
                   "Intestinal malabsorption (parasites, IBD, lymphangiectasia)",
                   "Hereditary VD-resistant rickets type II (rare, autosomal recessive – Pomeranians)"])

# ------------------------------------------------------------------
# 4. Clinical Signs
add_slide(prs,
          "Clinical Signs",
          bullets=["Age: 6 – 24 weeks (fast growth phase)",
                   "Lameness → reluctance to rise, exercise intolerance",
                   "Bone pain on palpation",
                   "Swollen metaphyses (wrists, hocks, stifles)",
                   "Bowed or angular limbs",
                   "Folding fractures of long bones & vertebrae",
                   "Stunted growth, loose teeth, alopecia in hereditary forms"])

# ------------------------------------------------------------------
# 5. Diagnosis – Clinicopathology
add_slide(prs,
          "Diagnosis – Clinicopathology",
          bullets=["Signalment + dietary history → high index of suspicion",
                   "Serum biochemistry",
                   "  – ↓ phosphorus (nutritional) or ↓ vitamin D",
                   "  – ↑ alkaline phosphatase (osteoblast activity)",
                   "  – ± mild hypocalcaemia (advanced)",
                   "Assay 25-OH-vitamin D – best reflection of body stores",
                   "Rule-out genetic forms via CYP27B1 / VDR gene tests"])

# ------------------------------------------------------------------
# 6. Diagnosis – Imaging
slide = add_slide(prs,
                  "Diagnosis – Imaging",
                  bullets=["Radiographs = gold standard (in vivo)",
                           "  – Generalised osteopenia",
                           "  – Widened, cupped, irregular growth plates",
                           "  – Flared metaphyses",
                           "  – Folding fractures",
                           "  – Angular limb deformity"],
                  picture="rickets_lateral.jpg",
                  notes="Lateral view of distal radius/ulna shows classic widening and cupping.")
dark_blue_bg(slide)

# ------------------------------------------------------------------
# 7. Radiographic Gallery
slide = add_slide(prs,
                  "Radiographic Gallery – Rickets",
                  bullets=["Compare with normal contralateral limb",
                           "Notice decreased radiopacity of cortices",
                           "Growth-plate width > 2 × normal",
                           "Secondary joint incongruency may lead to OA later"],
                  picture="rickets_ap.jpg",
                  notes="AP view of same dog – symmetrical physeal widening.")
dark_blue_bg(slide)

# ------------------------------------------------------------------
# 8. Treatment
add_slide(prs,
          "Treatment",
          bullets=["Correct the diet immediately",
                   "  – Balanced commercial puppy food (AAFCO growth)",
                   "  – Ca:P ratio 1.2–1.4 : 1; vitamin D ≥ 500 IU/1000 kcal",
                   "Sunlight exposure 30 min daily (UV-B 290–315 nm)",
                   "Specific supplementation if severe",
                   "  – Ca carbonate 50–100 mg kg⁻¹/day divided",
                   "  – Vitamin D3 (cholecalciferol) 1000–2000 IU/day × 4–6 wk",
                   "Analgesia (NSAIDs ± opioids) for pain/fractures",
                   "Orthopaedic intervention for pathological fractures"])

# ------------------------------------------------------------------
# 9. Prognosis & Management
add_slide(prs,
          "Prognosis & Management",
          bullets=["Excellent if diagnosed early (< 6 mo) & no irreversible physeal damage",
                   "Bone pain ↓ within 7–10 days; radiographic healing 4–6 weeks",
                   "Gradual resolution of angular deformity during remaining growth",
                   "Monitor every 2–4 weeks: body weight, gait, ALP, radiographs",
                   "Genetic cases require lifelong vitamin D analogues & Ca",
                   "Educate owners: avoid fad diets, feed growth-appropriate ration"])

# ------------------------------------------------------------------
# 10. Key Take-Home Messages
add_slide(prs,
          "Key Take-Home Messages",
          bullets=["Rickets is a preventable nutritional bone disease of puppies",
                   "Think of it in any lame, bow-legged, painful youngster",
                   "Radiographs give the diagnosis – look at the growth plates!",
                   "Treat the cause (diet + sunlight) – not just the bones",
                   "Prognosis is excellent with early intervention"])

# ------------------------------------------------------------------
# Save
# ------------------------------------------------------------------
prs.save("Rickets_in_the_Dog.pptx")
print("Presentation saved → Rickets_in_the_Dog.pptx")
  
