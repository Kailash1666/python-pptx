# build_rickets.py  –  build a detailed, image-rich PowerPoint on Rickets in the Dog
# Author: you  |  Date: 2025-11-21
import os, io, requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ---------- helper: download an image once ----------
def grab_img(url, file):
    if not os.path.exists(file):
        r = requests.get(url, timeout=30)
        with open(file, "wb") as f:
            f.write(r.content)
    return file

# ---------- radiology images (Wikimedia Commons – public domain) ----------
rad_lat = grab_img("https://upload.wikimedia.org/wikipedia/commons/3/3e/Rickets_lateral.jpg",
                   "rickets_lat.jpg")
rad_ap  = grab_img("https://upload.wikimedia.org/wikipedia/commons/2/2b/Rickets_AP.jpg",
                   "rickets_ap.jpg")

# ---------- slide builder ----------
def add_slide(prs, title, bullets=None, picture=None, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])   # blank layout
    # dark-blue background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)

    # title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.text, p.font.size, p.font.bold, p.font.color.rgb = title, Pt(34), True, RGBColor(255, 255, 255)

    if bullets:
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(7.5), Inches(5.5))
        tf = body.text_frame
        for line in bullets:
            pg = tf.add_paragraph()
            pg.text, pg.font.size, pg.font.color.rgb = line, Pt(20), RGBColor(255, 255, 255)

    if picture:
        slide.shapes.add_picture(picture, Inches(8.5), Inches(1.8), width=Inches(4.2))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

# ---------- build deck ----------
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)  # 16:9

# 1) Title
s0 = prs.slides.add_slide(prs.slide_layouts[0])
s0.shapes.title.text = "Rickets in the Dog"
s0.placeholders[1].text = "Etiology | Clinical Signs | Diagnosis | Treatment | Management\nVeterinary Continuing Education"
for para in (s0.shapes.title.text_frame.paragraphs[0],
             s0.placeholders[1].text_frame.paragraphs[0]):
    para.font.color.rgb = RGBColor(255, 255, 255)

# 2) Etiology
add_slide(prs,
          "Etiology & Risk Factors",
          bullets=["Nutritional (90 % cases)",
                   "  – All-meat diet → low Ca, low P, no vitamin D",
                   "  – Imbalanced homemade or raw diet",
                   "  – Excess Ca (≥3× normal) in giant breeds",
                   "Lack of sunlight → ↓cutaneous vitamin D₃",
                   "Intestinal malabsorption (parasites, IBD)",
                   "Hereditary vitamin-D-resistant forms (rare)",
                   "Poor milk intake or maternal malnutrition"],
          notes="Remember: rickets ONLY occurs in growing animals.")

# 3) Clinical Signs
add_slide(prs,
          "Clinical Signs",
          bullets=["Age 6 – 24 weeks; large breeds > small",
                   "Lameness / shifting-leg lameness",
                   "Pain on bone palpation; swollen metaphyses",
                   "Bowed fore-limbs, ‘knock-knees’",
                   "Folding fractures without trauma",
                   "Stiff gait, difficulty rising",
                   "Stunted growth, pot-belly, rachitic rosary",
                   "Depressed mood – pain induced"])

# 4) Diagnosis – clinicopathology
add_slide(prs,
          "Diagnosis – Clinicopathology",
          bullets=["Serum biochemistry",
                   "  – Low phosphorus (nutritional) or low vit-D",
                   "  – ↑ Alkaline phosphatase (osteoblast activity)",
                   "  – Ca normal or slightly ↓ (advanced disease)",
                   "Measure 25-OH-vitamin D – best body store indicator",
                   "Rule-out renal phosphate wasting (FeP) if suspected",
                   "Genetic tests for hereditary VD-resistant forms"])

# 5) Diagnosis – imaging 1
add_slide(prs,
          "Diagnosis – Imaging (1)",
          bullets=["Radiographs = gold standard",
                   "  – Generalised osteopenia",
                   "  – Widened, cupped, irregular growth plates",
                   "  – Flared metaphyses",
                   "  – Folding fractures / bowing",
                   "  – Thin cortices"],
          picture=rad_lat,
          notes="Lateral distal radius/ulna – classic physeal widening.")

# 6) Diagnosis – imaging 2
add_slide(prs,
          "Diagnosis – Imaging (2)",
          bullets=["Compare with normal littermate if possible",
                   "Growth-plate width ≥ 2 × normal",
                   "Irregular ‘moth-eaten’ metaphyseal border",
                   "Bone opacity ↓ 30-40 %",
                   "Secondary joint incongruency → future OA"],
          picture=rad_ap,
          notes="AP view – symmetrical physeal changes.")

# 7) Treatment
add_slide(prs,
          "Treatment",
          bullets=["Switch to balanced growth diet (AAFCO) immediately",
                   "Ca : P ratio 1.2–1.4 : 1; vitamin D ≥ 500 IU/1000 kcal",
                   "Sunlight 30 min daily (UV-B 290-315 nm)",
                   "Severe cases: supplement",
                   "  – Ca carbonate 50–100 mg/kg/day divided",
                   "  – Vit D₃ 1 000–2 000 IU/day × 4-6 wk",
                   "Analgesia (NSAIDs ± opioids)",
                   "Orthopaedic fixation for pathological fractures"])

# 8) Prognosis & Long-term Management
add_slide(prs,
          "Prognosis & Long-term Management",
          bullets=["Excellent if treated early (< 6 mo) & no physeal damage",
                   "Pain ↓ within 7-10 days; radiographic healing 4-6 wk",
                   "Angular deformity often self-corrects with remaining growth",
                   "Re-check q 2–4 wk: weight, gait, ALP, radiographs",
                   "Genetic cases need life-long vit-D analogues",
                   "Educate owners – avoid fad diets, feed growth ration"])

# 9) Key Take-Home Messages
add_slide(prs,
          "Key Take-Home Messages",
          bullets=["Rickets is a preventable nutritional bone disease of puppies",
                   "Lameness + bowed legs in a youngster → think rickets",
                   "Radiographs give the diagnosis – look at the physes!",
                   "Treat the cause (diet + sunlight) – not just the bones",
                   "Prognosis excellent with early correction"])

# ---------- save ----------
prs.save("Rickets_in_the_Dog.pptx")
print("Done → Rickets_in_the_Dog.pptx  (ready to copy to phone)")
